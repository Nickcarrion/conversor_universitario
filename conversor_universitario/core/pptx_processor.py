import gc
import os
import sys
import zipfile
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from utils.ocr_engine import ejecutar_ocr

def _abrir_presentacion(path: str) -> Presentation:
    """Abre PPTX y PPSX.

    PPSX tiene content-type 'slideshow' que python-pptx rechaza.
    Lo parcheamos reescribiendo [Content_Types].xml dentro del ZIP en memoria.
    """
    ext = os.path.splitext(path)[1].lower()
    if ext != ".ppsx":
        return Presentation(path)

    buf_out = BytesIO()
    with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(buf_out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(
                    b"presentationml.slideshow.main+xml",
                    b"presentationml.presentation.main+xml",
                )
            zout.writestr(item, data)

    buf_out.seek(0)
    return Presentation(buf_out)


def _texto_recursivo(shape) -> str:
    """Extrae texto de shapes normales y grupos (SmartArt, diagramas)."""
    texto = ""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                texto += _texto_recursivo(s)
        elif shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                linea = para.text.strip()
                if linea:
                    texto += linea + "\n"
    except Exception:
        pass
    return texto


def _procesar_imagen_shape(shape, slide_idx: int, img_dir: str, lang: str, img_prefix: str = "") -> str:
    """Guarda la imagen del shape, aplica OCR y retorna el bloque Markdown."""
    try:
        image = shape.image
        img_bytes = image.blob
        img_ext = image.ext or "png"
        filename = f"{img_prefix}slide_{slide_idx + 1}_img_{shape.shape_id}.{img_ext}"
        img_path = os.path.join(img_dir, filename)

        os.makedirs(img_dir, exist_ok=True)
        with open(img_path, "wb") as f:
            f.write(img_bytes)

        with Image.open(BytesIO(img_bytes)) as img_pil:
            txt_ocr = ejecutar_ocr(img_pil, lang)

        rel_path = img_path.replace("\\", "/")
        bloque = f"\n![imagen diapositiva]({rel_path})\n"
        if txt_ocr and not txt_ocr.startswith("["):
            bloque += f"\n> **OCR de imagen:** {txt_ocr.replace(chr(10), ' ')}\n"
        return bloque
    except Exception as e:
        print(f"  [Advertencia] No se pudo procesar imagen en slide {slide_idx + 1}: {e}", file=sys.stderr)
        return ""


def convert_pptx(path: str, lang: str = "spa", img_dir: str = None, img_prefix: str = "") -> str:
    base_name = os.path.splitext(os.path.basename(path))[0]
    if img_dir is None:
        img_dir = os.path.join(os.path.dirname(path), f"{base_name}_images")

    prs = _abrir_presentacion(path)
    total = len(prs.slides)
    partes = []

    for idx, slide in enumerate(prs.slides):
        titulo = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                titulo = f" — {slide.shapes.title.text.strip()}"
        except Exception:
            pass

        partes.append(f"\n## Diapositiva {idx + 1}{titulo}\n\n")

        texto_slide = ""
        imagenes_md = ""

        for shape in slide.shapes:
            texto_slide += _texto_recursivo(shape)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                imagenes_md += _procesar_imagen_shape(shape, idx, img_dir, lang, img_prefix)

        if texto_slide.strip():
            partes.append(texto_slide)
        if imagenes_md:
            partes.append(imagenes_md)

        # Notas del presentador
        try:
            if slide.has_notes_slide:
                notas_tf = slide.notes_slide.notes_text_frame
                if notas_tf:
                    notas = notas_tf.text.strip()
                    if notas:
                        partes.append(f"\n> **Notas del Presentador:**\n> {notas.replace(chr(10), ' ')}\n")
        except Exception:
            pass

        partes.append("\n---\n")

        # Liberar memoria tras cada slide
        gc.collect()
        print(f"  [{idx + 1}/{total}] Diapositiva procesada", end="\r")

    print()
    return "".join(partes)
